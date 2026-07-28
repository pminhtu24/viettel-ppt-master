"""Core PPTX assembly: create_pptx_with_native_svg."""

from __future__ import annotations

import hashlib
import mimetypes
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Any

from pptx import Presentation

from .drawingml_converter import convert_svg_to_slide_shapes
from .pptx_dimensions import (
    CANVAS_FORMATS,
    get_slide_dimensions, get_pixel_dimensions,
    get_viewbox_dimensions, detect_format_from_svg,
)
from .pptx_slide_xml import (
    ANIMATIONS_AVAILABLE, TRANSITIONS,
)

# Re-import create_transition_xml only if available
try:
    from pptx_animations import (
        create_transition_xml,
        create_sequence_timing_xml,
        pick_animation_effect,
    )
except ImportError:
    create_transition_xml = None
    create_sequence_timing_xml = None
    pick_animation_effect = None


_IMAGE_CONTENT_TYPES = {
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'gif': 'image/gif',
    'webp': 'image/webp',
    'svg': 'image/svg+xml',
    'bmp': 'image/bmp',
    'emf': 'image/x-emf',
    'tif': 'image/tiff',
    'tiff': 'image/tiff',
    'wmf': 'image/x-wmf',
}


def _content_type_for_extension(ext: str) -> str:
    clean = ext.lower().lstrip('.')
    content_type = _IMAGE_CONTENT_TYPES.get(clean) or mimetypes.guess_type(f'x.{clean}')[0]
    if not content_type:
        raise ValueError(f"Unknown media content type for extension: {ext}")
    return content_type


def _as_dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _to_float(value: Any, default: float) -> float:
    if value is None:
        return default
    try:
        number = float(value)
    except (TypeError, ValueError):
        return default
    return number if number >= 0 else default


def _slide_config(animation_config: dict[str, Any] | None, svg_stem: str) -> dict[str, Any]:
    if not animation_config:
        return {}
    slides = _as_dict(animation_config.get('slides'))
    return _as_dict(slides.get(svg_stem))


def _slide_transition_settings(
    slide_cfg: dict[str, Any],
    transition: str | None,
    duration: float,
    auto_advance: float | None,
    cli_overrides: dict[str, bool],
) -> tuple[str | None, float, float | None]:
    trans_cfg = _as_dict(slide_cfg.get('transition'))
    effect = transition
    if not cli_overrides.get('transition') and 'effect' in trans_cfg:
        cfg_effect = str(trans_cfg.get('effect'))
        effect = None if cfg_effect == 'none' else cfg_effect
    if not cli_overrides.get('transition_duration'):
        duration = _to_float(trans_cfg.get('duration'), duration)
    if not cli_overrides.get('auto_advance') and 'auto_advance' in trans_cfg:
        auto_advance = _to_float(trans_cfg.get('auto_advance'), auto_advance or 0)
    return effect, duration, auto_advance


def _slide_animation_settings(
    slide_cfg: dict[str, Any],
    animation: str | None,
    duration: float,
    stagger: float,
    trigger: str,
    cli_overrides: dict[str, bool],
) -> tuple[str | None, float, float, str]:
    anim_cfg = _as_dict(slide_cfg.get('animation'))
    effect = animation
    if not cli_overrides.get('animation') and 'effect' in anim_cfg:
        cfg_effect = str(anim_cfg.get('effect'))
        effect = None if cfg_effect == 'none' else cfg_effect
    if not cli_overrides.get('animation_duration'):
        duration = _to_float(anim_cfg.get('duration'), duration)
    if not cli_overrides.get('animation_stagger'):
        stagger = _to_float(anim_cfg.get('stagger'), stagger)
    if not cli_overrides.get('animation_trigger') and anim_cfg.get('trigger'):
        trigger = str(anim_cfg.get('trigger'))
    return effect, duration, stagger, trigger


def _build_sequence_targets(
    anim_targets: list[tuple[int, str]],
    slide_cfg: dict[str, Any],
    animation: str,
    duration: float,
    stagger: float,
    mixed_animation_offset: int,
) -> tuple[list[tuple[int, int, str, float]], int]:
    groups_cfg = _as_dict(slide_cfg.get('groups'))
    ordered: list[tuple[int, int, int, str, dict[str, Any]]] = []
    for idx, (sid, svg_id) in enumerate(anim_targets):
        group_cfg = _as_dict(groups_cfg.get(svg_id))
        if str(group_cfg.get('effect', '')).lower() == 'none':
            continue
        order_value = group_cfg.get('order')
        try:
            order = int(order_value)
            has_order = 0
        except (TypeError, ValueError):
            order = idx
            has_order = 1
        group_entry = dict(group_cfg)
        group_entry['_shape_id'] = sid
        ordered.append((has_order, order, idx, svg_id, group_entry))

    ordered.sort(key=lambda item: (item[0], item[1], item[2]))

    seq_targets: list[tuple[int, int, str, float]] = []
    for seq_idx, (_has_order, _order, _original_idx, _svg_id, group_cfg) in enumerate(ordered):
        shape_id = int(group_cfg['_shape_id'])
        raw_effect = group_cfg.get('effect')
        if raw_effect in ('mixed', 'random'):
            effect = pick_animation_effect(str(raw_effect), seq_idx, mixed_animation_offset)
        else:
            effect = str(raw_effect or pick_animation_effect(
                animation, seq_idx, mixed_animation_offset,
            ))
        item_duration = _to_float(group_cfg.get('duration'), duration)
        delay_seconds = _to_float(
            group_cfg.get('delay'),
            0 if seq_idx == 0 else stagger,
        )
        seq_targets.append((shape_id, int(delay_seconds * 1000), effect, item_duration))

    mixed_count = 0
    if animation == 'mixed':
        mixed_count = sum(1 for _target in seq_targets[1:])
    return seq_targets, mixed_count


def create_pptx_with_native_svg(
    svg_files: list[Path],
    output_path: Path,
    canvas_format: str | None = None,
    verbose: bool = True,
    transition: str | None = 'fade',
    transition_duration: float = 0.5,
    auto_advance: float | None = None,
    animation: str | None = None,
    animation_duration: float = 0.4,
    animation_stagger: float = 0.5,
    animation_trigger: str = 'after-previous',
    animation_config: dict[str, Any] | None = None,
    animation_cli_overrides: dict[str, bool] | None = None,
) -> bool:
    """Create a PPTX file with native SVG.

    Args:
        svg_files: List of SVG files.
        output_path: Output PPTX path.
        canvas_format: Canvas format key.
        verbose: Whether to output detailed information.
        transition: Transition effect name.
        transition_duration: Transition duration in seconds.
        auto_advance: Auto-advance interval in seconds.
        animation: Per-element entrance animation mode (single effect name,
            'mixed', 'random', or None to disable). Native shapes mode only.
        animation_duration: Per-element entrance duration in seconds.
        animation_stagger: Delay between elements in ``after-previous``
            trigger mode (seconds). Ignored otherwise.
        animation_trigger: PowerPoint Start mode — ``'after-previous'`` (default),
            ``'on-click'``, or ``'with-previous'``.
        animation_config: Optional sidecar overrides loaded from animations.json.
        animation_cli_overrides: Flags indicating explicit CLI overrides.
    Returns:
        Whether all slides were successfully created.
    """
    if not svg_files:
        print("Error: No SVG files found")
        return False

    # Auto-detect canvas format or get dimensions from viewBox
    custom_pixels: tuple[int, int] | None = None
    if canvas_format is None:
        canvas_format = detect_format_from_svg(svg_files[0])
        if canvas_format and verbose:
            format_name = CANVAS_FORMATS.get(canvas_format, {}).get('name', canvas_format)
            print(f"  Detected canvas format: {format_name}")

    if canvas_format is None:
        custom_pixels = get_viewbox_dimensions(svg_files[0])
        if custom_pixels and verbose:
            print(f"  Using SVG viewBox dimensions: {custom_pixels[0]} x {custom_pixels[1]} px")

    if canvas_format is None and custom_pixels is None:
        canvas_format = 'ppt169'
        if verbose:
            print(f"  Using default format: PPT 16:9")

    width_emu, height_emu = get_slide_dimensions(canvas_format or 'ppt169', custom_pixels)
    pixel_width, pixel_height = get_pixel_dimensions(canvas_format or 'ppt169', custom_pixels)

    if verbose:
        print(f"  Slide dimensions: {pixel_width} x {pixel_height} px")
        print(f"  SVG file count: {len(svg_files)}")
        print("  Mode: Native DrawingML shapes (directly editable)")
        if transition:
            trans_name = TRANSITIONS.get(transition, {}).get('name', transition) if TRANSITIONS else transition
            print(f"  Transition effect: {trans_name}")
        print()

    animation_cli_overrides = animation_cli_overrides or {}

    temp_dir = Path(tempfile.mkdtemp())

    try:
        # Create base PPTX with python-pptx
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu

        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)

        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))

        # Extract PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)

        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)

        success_count = 0
        media_cache: dict[tuple[str, str], str] = {}
        image_exts_used: set[str] = set()
        mixed_animation_offset = 0

        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i

            try:
                slide_cfg = _slide_config(animation_config, svg_path.stem)
                slide_xml, media_files_dict, rel_entries, anim_targets = (
                    convert_svg_to_slide_shapes(
                        svg_path, slide_num=slide_num, verbose=verbose,
                    )
                )
                slide_transition, slide_transition_duration, slide_auto_advance = (
                    _slide_transition_settings(
                        slide_cfg,
                        transition,
                        transition_duration,
                        auto_advance,
                        animation_cli_overrides,
                    )
                )
                (
                    slide_animation,
                    slide_animation_duration,
                    slide_animation_stagger,
                    slide_animation_trigger,
                ) = _slide_animation_settings(
                    slide_cfg,
                    animation,
                    animation_duration,
                    animation_stagger,
                    animation_trigger,
                    animation_cli_overrides,
                )

                # OOXML requires transition before timing inside p:sld.
                if slide_transition and ANIMATIONS_AVAILABLE and create_transition_xml:
                    transition_xml = '\n' + create_transition_xml(
                        effect=slide_transition,
                        duration=slide_transition_duration,
                        advance_after=slide_auto_advance,
                    )
                    slide_xml = slide_xml.replace(
                        '</p:sld>',
                        transition_xml + '\n</p:sld>',
                    )

                if (slide_animation and slide_animation != 'none'
                        and create_sequence_timing_xml
                        and pick_animation_effect
                        and anim_targets):
                    seq_targets, mixed_count = _build_sequence_targets(
                        anim_targets,
                        slide_cfg,
                        slide_animation,
                        slide_animation_duration,
                        slide_animation_stagger,
                        mixed_animation_offset,
                    )
                    if slide_animation == 'mixed':
                        mixed_animation_offset += mixed_count
                    timing_xml = '\n' + create_sequence_timing_xml(
                        seq_targets, duration=slide_animation_duration,
                        trigger=slide_animation_trigger,
                    )
                    slide_xml = slide_xml.replace(
                        '</p:sld>',
                        timing_xml + '\n</p:sld>',
                    )

                slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                with open(slide_xml_path, 'w', encoding='utf-8') as f:
                    f.write(slide_xml)

                media_name_map: dict[str, str] = {}
                for media_name, media_data in media_files_dict.items():
                    ext = media_name.rsplit('.', 1)[-1].lower()
                    media_hash = hashlib.sha256(media_data).hexdigest()
                    cache_key = (ext, media_hash)
                    cached_name = media_cache.get(cache_key)

                    if cached_name is None:
                        cached_name = f'image_{media_hash[:16]}.{ext}'
                        media_cache[cache_key] = cached_name
                        with open(media_dir / cached_name, 'wb') as f:
                            f.write(media_data)

                    media_name_map[media_name] = cached_name

                for rel in rel_entries:
                    target = rel.get('target', '')
                    if not target.startswith('../media/'):
                        continue
                    media_name = target.split('../media/', 1)[1]
                    mapped_name = media_name_map.get(media_name)
                    if mapped_name:
                        rel['target'] = f'../media/{mapped_name}'

                rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                rels_dir.mkdir(exist_ok=True)
                rels_path = rels_dir / f'slide{slide_num}.xml.rels'

                extra_rels = ''
                for rel in rel_entries:
                    extra_rels += (
                        f'\n  <Relationship Id="{rel["id"]}" '
                        f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
                    )

                rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>{extra_rels}
</Relationships>'''
                with open(rels_path, 'w', encoding='utf-8') as f:
                    f.write(rels_xml)

                for media_name in media_name_map.values():
                    ext = media_name.rsplit('.', 1)[-1].lower()
                    _content_type_for_extension(ext)
                    image_exts_used.add(ext)

                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} (Native)")

                success_count += 1

            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - Error: {e}")
                raise

        # Update [Content_Types].xml
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()

        types_to_add: list[str] = []
        for ext in sorted(image_exts_used):
            if f'Extension="{ext}"' not in content_types:
                types_to_add.append(
                    f'  <Default Extension="{ext}" ContentType="{_content_type_for_extension(ext)}"/>'
                )

        if types_to_add:
            content_types = content_types.replace(
                '</Types>', '\n'.join(types_to_add) + '\n</Types>',
            )
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Repackage PPTX to a temporary file first. The public output path is
        # replaced only after every slide and relationship has succeeded.
        temp_output_path = temp_dir / 'result.pptx'
        with zipfile.ZipFile(temp_output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)
        shutil.move(str(temp_output_path), str(output_path))

        if verbose:
            print()
            print(f"[Done] Saved: {output_path}")
            print(f"  Succeeded: {success_count}, Failed: {len(svg_files) - success_count}")

        return success_count == len(svg_files)

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
